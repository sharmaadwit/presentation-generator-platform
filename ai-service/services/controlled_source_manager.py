import asyncio
import asyncpg
import os
from typing import List, Dict, Any, Optional
import logging
from pptx import Presentation
import json
import boto3
import tempfile
from urllib.parse import urlparse
import base64

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ControlledSourceManager:
    def __init__(self):
        self.database_url = os.getenv('DATABASE_URL', 'postgresql://localhost:5432/presentation_generator')
        self.connection_pool = None
        self.s3_client = None
        self._init_s3_client()
    
    def _init_s3_client(self):
        """Initialize S3 client with AWS credentials"""
        try:
            self.s3_client = boto3.client(
                's3',
                region_name=os.getenv('AWS_REGION', 'ap-south-1')
            )
            logger.info("S3 client initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize S3 client: {e}")
            self.s3_client = None
    
    async def connect(self):
        """Connect to the database"""
        try:
            self.connection_pool = await asyncpg.create_pool(
                self.database_url,
                min_size=1,
                max_size=10,
                command_timeout=60
            )
            logger.info("Connected to controlled source database")
        except Exception as e:
            logger.error(f"Failed to connect to database: {e}")
            raise e
    
    async def close(self):
        """Close database connection"""
        if self.connection_pool:
            await self.connection_pool.close()
    
    async def get_approved_sources(self, industry: str = None, tags: List[str] = None, limit: int = 20) -> List[Dict[str, Any]]:
        """Get approved presentation sources with optional filtering"""
        try:
            if not self.connection_pool:
                await self.connect()
                
            async with self.connection_pool.acquire() as conn:
                # Build query with optional filters
                where_conditions = ["ps.status = 'approved'"]
                params = []
                param_count = 0
                
                if industry:
                    param_count += 1
                    where_conditions.append(f"ps.industry = ${param_count}")
                    params.append(industry)
                
                if tags:
                    param_count += 1
                    where_conditions.append(f"ps.tags && ${param_count}")
                    params.append(tags)
                
                where_clause = " AND ".join(where_conditions)
                params.append(limit)
                
                query = f"""
                    SELECT ps.*, 
                           COUNT(ss.id) as slide_count
                    FROM presentation_sources ps
                    LEFT JOIN source_slides ss ON ps.id = ss.source_id
                    WHERE {where_clause}
                    GROUP BY ps.id
                    ORDER BY ps.relevance_score DESC, ps.created_at DESC
                    LIMIT ${param_count + 1}
                """
                
                rows = await conn.fetch(query, *params)
                return [dict(row) for row in rows]
                
        except Exception as e:
            logger.error(f"Error getting approved sources: {e}")
            return []

    async def get_approved_sources_for_industry(self, industry: str) -> List[Dict[str, Any]]:
        """Get all approved presentation sources for a specific industry"""
        try:
            async with self.connection_pool.acquire() as conn:
                rows = await conn.fetch("""
                    SELECT ps.*, 
                           COUNT(ss.id) as slide_count
                    FROM presentation_sources ps
                    LEFT JOIN source_slides ss ON ps.id = ss.source_id
                    WHERE ps.status = 'approved' 
                    AND ps.industry = $1
                    GROUP BY ps.id
                    ORDER BY ps.relevance_score DESC, ps.created_at DESC
                """, industry)
                
                return [dict(row) for row in rows]
        except Exception as e:
            logger.error(f"Error getting approved sources: {e}")
            return []
    
    async def get_all_approved_sources(self) -> List[Dict[str, Any]]:
        """Get all approved presentation sources"""
        try:
            async with self.connection_pool.acquire() as conn:
                rows = await conn.fetch("""
                    SELECT ps.*, 
                           COUNT(ss.id) as slide_count
                    FROM presentation_sources ps
                    LEFT JOIN source_slides ss ON ps.id = ss.source_id
                    WHERE ps.status = 'approved'
                    GROUP BY ps.id
                    ORDER BY ps.relevance_score DESC, ps.created_at DESC
                """)
                
                return [dict(row) for row in rows]
        except Exception as e:
            logger.error(f"Error getting all approved sources: {e}")
            return []
    
    async def get_source_slides(self, source_id: str) -> List[Dict[str, Any]]:
        """Get all slides from a specific approved source"""
        try:
            async with self.connection_pool.acquire() as conn:
                rows = await conn.fetch("""
                    SELECT * FROM source_slides 
                    WHERE source_id = $1 
                    ORDER BY slide_index
                """, source_id)
                
                return [dict(row) for row in rows]
        except Exception as e:
            logger.error(f"Error getting source slides: {e}")
            return []
    
    async def extract_slides_from_source(self, source_id: str, file_path: str, industry: str = None) -> List[Dict[str, Any]]:
        """Extract slides from a presentation file and store them in the database"""
        temp_file = None
        try:
            # Handle S3 files
            if file_path.startswith('s3://'):
                temp_file = await self._download_s3_file(file_path)
                if not temp_file:
                    logger.error(f"Failed to download S3 file: {file_path}")
                    return []
                presentation_path = temp_file
            else:
                presentation_path = file_path
            
            # Load the presentation
            prs = Presentation(presentation_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_data = self._extract_slide_data(slide, i)
                slide_data['source_id'] = source_id
                slides.append(slide_data)
            
            # Store slides in database with industry information
            await self._store_slides(slides, industry)
            
            logger.info(f"Extracted {len(slides)} slides from source {source_id} for industry {industry}")
            return slides
            
        except Exception as e:
            logger.error(f"Error extracting slides from source {source_id}: {e}")
            return []
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.unlink(temp_file)
                except Exception as e:
                    logger.warning(f"Failed to clean up temp file {temp_file}: {e}")
    
    async def _download_s3_file(self, s3_path: str) -> Optional[str]:
        """Download file from S3 to temporary location"""
        if not self.s3_client:
            logger.error("S3 client not initialized")
            return None
        
        try:
            # Parse S3 URL
            parsed = urlparse(s3_path)
            bucket = parsed.netloc
            key = parsed.path.lstrip('/')
            
            # Create temporary file
            temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
            os.close(temp_fd)
            
            # Download file from S3
            logger.info(f"Downloading S3 file: s3://{bucket}/{key}")
            self.s3_client.download_file(bucket, key, temp_path)
            
            logger.info(f"Successfully downloaded S3 file to: {temp_path}")
            return temp_path
            
        except Exception as e:
            logger.error(f"Error downloading S3 file {s3_path}: {e}")
            return None
    
    def _extract_slide_data(self, slide, index: int) -> Dict[str, Any]:
        """Extract data from a single slide with enhanced visual element preservation"""
        slide_data = {
            'slide_index': index,
            'title': '',
            'content': '',
            'image_url': None,
            'images': [],  # Store multiple images
            'slide_type': 'content',
            'metadata': {},
            'formatting': {},  # Store formatting information
            'layout_info': {}  # Store layout information
        }
        
        try:
            # Extract title
            if slide.shapes.title:
                slide_data['title'] = slide.shapes.title.text
                # Extract title formatting
                slide_data['formatting']['title'] = self._extract_text_formatting(slide.shapes.title)
            
            # Extract content from text boxes with formatting
            content_parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if shape != slide.shapes.title:  # Don't duplicate title
                        content_parts.append(shape.text)
                        # Extract formatting for each text shape
                        shape_formatting = self._extract_text_formatting(shape)
                        if shape_formatting:
                            slide_data['formatting'][f'text_shape_{len(content_parts)}'] = shape_formatting
            
            slide_data['content'] = '\n'.join(content_parts)
            
            # Extract images and visual elements
            slide_data['images'] = self._extract_images(slide)
            
            # Extract slide background and layout information
            slide_data['layout_info'] = self._extract_layout_info(slide)
            
            # Determine slide type
            slide_data['slide_type'] = self._classify_slide_type(slide_data)
            
            # Enhanced metadata
            slide_data['metadata'] = {
                'shapes_count': len(slide.shapes),
                'has_images': len(slide_data['images']) > 0,
                'image_count': len(slide_data['images']),
                'text_shapes': len([s for s in slide.shapes if hasattr(s, 'text')]),
                'has_background': hasattr(slide.background, 'fill'),
                'slide_layout': getattr(slide.slide_layout, 'name', 'unknown') if hasattr(slide, 'slide_layout') else 'unknown'
            }
            
        except Exception as e:
            logger.error(f"Error extracting slide data: {e}")
        
        return slide_data
    
    def _extract_text_formatting(self, shape) -> Dict[str, Any]:
        """Extract text formatting information from a shape"""
        formatting = {}
        
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame = shape.text_frame
                
                # Extract paragraph formatting
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    formatting['paragraph'] = {
                        'alignment': str(para.alignment) if para.alignment else None,
                        'space_before': para.space_before,
                        'space_after': para.space_after
                    }
                    
                    # Extract font formatting
                    if para.runs:
                        run = para.runs[0]
                        # Handle color safely
                        color_value = None
                        if run.font.color:
                            try:
                                if hasattr(run.font.color, 'rgb'):
                                    color_value = str(run.font.color.rgb)
                                elif hasattr(run.font.color, 'theme_color'):
                                    color_value = f"theme:{run.font.color.theme_color}"
                                else:
                                    color_value = str(run.font.color)
                            except Exception:
                                color_value = None
                        
                        formatting['font'] = {
                            'name': run.font.name,
                            'size': run.font.size.pt if run.font.size else None,
                            'bold': run.font.bold,
                            'italic': run.font.italic,
                            'underline': run.font.underline,
                            'color': color_value
                        }
        except Exception as e:
            logger.warning(f"Could not extract text formatting: {e}")
        
        return formatting
    
    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """Extract images from slide including GIFs and other formats"""
        images = []
        
        try:
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image'):
                    try:
                        # Get image data
                        image_blob = shape.image.blob if hasattr(shape.image, 'blob') else None
                        
                        # Determine image format
                        content_type = shape.image.content_type if hasattr(shape.image, 'content_type') else 'unknown'
                        
                        # Check if it's a GIF
                        is_gif = False
                        if content_type and 'gif' in content_type.lower():
                            is_gif = True
                        elif image_blob and image_blob.startswith(b'GIF'):
                            is_gif = True
                            content_type = 'image/gif'
                        
                        image_data = {
                            'index': i,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height,
                            'image_format': content_type,
                            'is_gif': is_gif,
                            'image_data': base64.b64encode(image_blob).decode('utf-8') if image_blob else None
                            # Note: Removed image_blob to avoid JSON serialization issues
                            # Raw blob data is not needed for database storage
                        }
                        
                        # Try to extract image filename or identifier
                        if hasattr(shape.image, 'filename'):
                            image_data['filename'] = shape.image.filename
                        
                        # Log the image type
                        if is_gif:
                            logger.info(f"Extracted GIF image: {image_data.get('filename', f'image_{i}')}")
                        else:
                            logger.info(f"Extracted image: {image_data.get('filename', f'image_{i}')} ({content_type})")
                        
                        images.append(image_data)
                        
                    except Exception as e:
                        logger.warning(f"Could not extract image {i}: {e}")
                        continue
                        
        except Exception as e:
            logger.warning(f"Could not extract images from slide: {e}")
        
        return images
    
    def _extract_layout_info(self, slide) -> Dict[str, Any]:
        """Extract layout and background information"""
        layout_info = {}
        
        try:
            # Extract background information
            if hasattr(slide, 'background') and slide.background:
                background = slide.background
                layout_info['background'] = {
                    'has_fill': hasattr(background, 'fill'),
                    'fill_type': str(background.fill.type) if hasattr(background, 'fill') else None
                }
            
            # Extract slide dimensions
            layout_info['dimensions'] = {
                'width': slide.slide_layout.slide_width if hasattr(slide.slide_layout, 'slide_width') else None,
                'height': slide.slide_layout.slide_height if hasattr(slide.slide_layout, 'slide_height') else None
            }
            
            # Extract slide layout information
            if hasattr(slide, 'slide_layout'):
                layout_info['layout'] = {
                    'name': getattr(slide.slide_layout, 'name', 'unknown'),
                    'layout_id': getattr(slide.slide_layout, 'layout_id', None)
                }
                
        except Exception as e:
            logger.warning(f"Could not extract layout info: {e}")
        
        return layout_info
    
    def _classify_slide_type(self, slide_data: Dict[str, Any]) -> str:
        """Classify the type of slide based on content"""
        title = slide_data.get('title', '').lower()
        content = slide_data.get('content', '').lower()
        
        if any(word in title for word in ['title', 'agenda', 'overview']):
            return 'title'
        elif any(word in content for word in ['chart', 'graph', 'data', 'statistics']):
            return 'chart'
        elif any(word in content for word in ['quote', 'testimonial', 'feedback']):
            return 'quote'
        elif any(word in content for word in ['conclusion', 'summary', 'next steps']):
            return 'conclusion'
        else:
            return 'content'
    
    async def _store_slides(self, slides: List[Dict[str, Any]], industry: str = None):
        """Store extracted slides in the database with enhanced visual data"""
        try:
            async with self.connection_pool.acquire() as conn:
                for slide in slides:
                    await conn.execute("""
                        INSERT INTO source_slides (
                            source_id, slide_index, title, content, 
                            image_url, slide_type, metadata, 
                            images, formatting, layout_info, industry
                        ) VALUES ($1, $2, $3, $4, $5, $6, $7, $8, $9, $10, $11)
                    """, 
                    slide['source_id'],
                    slide['slide_index'],
                    slide['title'],
                    slide['content'],
                    slide['image_url'],
                    slide['slide_type'],
                    json.dumps(slide['metadata']),
                    json.dumps(slide.get('images', [])),
                    json.dumps(slide.get('formatting', {})),
                    json.dumps(slide.get('layout_info', {})),
                    industry
                    )
        except Exception as e:
            logger.error(f"Error storing slides: {e}")
            raise e
    
    async def store_presentation_source(
        self, 
        source_id: str, 
        title: str, 
        description: str, 
        industry: str, 
        tags: list, 
        file_path: str, 
        mime_type: str
    ):
        """Store presentation source information in the database"""
        try:
            async with self.connection_pool.acquire() as conn:
                await conn.execute("""
                    INSERT INTO presentation_sources (
                        id, title, description, industry, tags, 
                        file_path, source_type, status, 
                        uploaded_by, created_at, updated_at, metadata
                    ) VALUES ($1, $2, $3, $4, $5, $6, $7, $8, $9, CURRENT_TIMESTAMP, CURRENT_TIMESTAMP, $10)
                    ON CONFLICT (id) DO UPDATE SET
                        title = EXCLUDED.title,
                        description = EXCLUDED.description,
                        industry = EXCLUDED.industry,
                        tags = EXCLUDED.tags,
                        file_path = EXCLUDED.file_path,
                        metadata = EXCLUDED.metadata,
                        updated_at = CURRENT_TIMESTAMP
                """, 
                source_id, title, description, industry, tags, 
                file_path, 'uploaded', 'approved', 'system',
                json.dumps({'mime_type': mime_type})
                )
                
                logger.info(f"Stored presentation source: {source_id} for industry: {industry}")
                
        except Exception as e:
            logger.error(f"Error storing presentation source: {e}")
            raise e

    async def search_slides_by_criteria(
        self, 
        industry: str, 
        use_case: str, 
        customer: str,
        target_audience: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """Search for relevant slides from approved sources based on criteria"""
        try:
            async with self.connection_pool.acquire() as conn:
                # First get approved sources for the industry
                sources = await conn.fetch("""
                    SELECT id FROM presentation_sources 
                    WHERE status = 'approved' AND industry = $1
                """, industry)
                
                if not sources:
                    logger.warning(f"No approved sources found for industry: {industry}")
                    return []
                
                source_ids = [row['id'] for row in sources]
                
                # Search for relevant slides
                query = f"""
                    SELECT ss.*, ps.title as source_title, ps.description as source_description
                    FROM source_slides ss
                    JOIN presentation_sources ps ON ss.source_id = ps.id
                    WHERE ss.source_id = ANY($1)
                    AND (
                        ss.title ILIKE $2 OR 
                        ss.content ILIKE $2 OR
                        ps.title ILIKE $2 OR
                        ps.description ILIKE $2
                    )
                    ORDER BY ps.relevance_score DESC, ss.slide_index
                """
                
                search_term = f"%{use_case}%{customer}%"
                if target_audience:
                    search_term += f"%{target_audience}%"
                
                rows = await conn.fetch(query, source_ids, search_term)
                
                return [dict(row) for row in rows]
                
        except Exception as e:
            logger.error(f"Error searching slides: {e}")
            return []
    
    async def get_source_statistics(self) -> Dict[str, Any]:
        """Get statistics about approved sources"""
        try:
            async with self.connection_pool.acquire() as conn:
                # Get overall stats
                stats = await conn.fetchrow("""
                    SELECT 
                        COUNT(DISTINCT ps.id) as total_sources,
                        COUNT(ss.id) as total_slides,
                        COUNT(DISTINCT ps.industry) as industries_covered
                    FROM presentation_sources ps
                    LEFT JOIN source_slides ss ON ps.id = ss.source_id
                    WHERE ps.status = 'approved'
                """)
                
                # Get industry breakdown
                industry_stats = await conn.fetch("""
                    SELECT 
                        ps.industry,
                        COUNT(DISTINCT ps.id) as source_count,
                        COUNT(ss.id) as slide_count
                    FROM presentation_sources ps
                    LEFT JOIN source_slides ss ON ps.id = ss.source_id
                    WHERE ps.status = 'approved'
                    GROUP BY ps.industry
                    ORDER BY slide_count DESC
                """)
                
                return {
                    'overview': dict(stats),
                    'industry_breakdown': [dict(row) for row in industry_stats]
                }
                
        except Exception as e:
            logger.error(f"Error getting source statistics: {e}")
            return {}
    
    async def validate_source_access(self, source_id: str) -> bool:
        """Validate that a source is approved and accessible"""
        try:
            async with self.connection_pool.acquire() as conn:
                row = await conn.fetchrow("""
                    SELECT status FROM presentation_sources 
                    WHERE id = $1
                """, source_id)
                
                return row and row['status'] == 'approved'
                
        except Exception as e:
            logger.error(f"Error validating source access: {e}")
            return False
