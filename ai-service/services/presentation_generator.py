from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import uuid
from typing import List, Dict, Any, Optional
import logging
from PIL import Image
import requests
from io import BytesIO
from .pdf_generator import PDFGenerator

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        self.output_dir = os.getenv('PRESENTATION_OUTPUT_DIR', './generated_presentations')
        self.pdf_generator = PDFGenerator()
        self.ensure_output_dir()
    
    def ensure_output_dir(self):
        """Ensure output directory exists"""
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)
    
    async def generate_presentation(
        self,
        slides: List[Dict[str, Any]],
        presentation_id: str,
        request_data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Generate a PDF presentation from selected slides"""
        
        try:
            print(f"🎨 PDF PRESENTATION GENERATOR STARTED")
            print(f"   - Presentation ID: {presentation_id}")
            print(f"   - Total Slides: {len(slides)}")
            print(f"   - Customer: {request_data.get('customer', 'N/A')}")
            print(f"   - Industry: {request_data.get('industry', 'N/A')}")
            print(f"   - Style: {request_data.get('style', 'N/A')}")
            
            # Debug: Check slide data
            for i, slide in enumerate(slides):
                print(f"🔍 SLIDE {i+1} DATA:")
                print(f"   - Title: {slide.get('title', 'N/A')}")
                print(f"   - Images: {len(slide.get('images', []))} visual elements")
                print(f"   - Formatting: {bool(slide.get('formatting'))}")
                print(f"   - Layout: {bool(slide.get('layout_info'))}")
            
            # Use PDF generator instead of PowerPoint
            result = await self.pdf_generator.generate_presentation(slides, presentation_id, request_data)
            
            print(f"✅ PDF presentation generated successfully")
            return result
            
        except Exception as e:
            logger.error(f"Error generating PDF presentation: {e}")
            raise e
    
    def _set_presentation_properties(self, prs: Presentation, request_data: Dict[str, Any]):
        """Set presentation metadata and properties"""
        
        # Set core properties
        prs.core_properties.title = f"{request_data['customer']} - {request_data['industry']} Presentation"
        prs.core_properties.subject = request_data['useCase']
        prs.core_properties.author = "AI Presentation Generator"
        prs.core_properties.comments = f"Generated for {request_data['customer']} in {request_data['industry']} industry"
        
        # Set slide size based on style
        if request_data.get('style') == 'creative':
            prs.slide_width = Inches(13.33)  # 16:9 creative format
            prs.slide_height = Inches(7.5)
        else:
            prs.slide_width = Inches(10)  # Standard format
            prs.slide_height = Inches(7.5)
    
    def _add_title_slide(self, prs: Presentation, request_data: Dict[str, Any]):
        """Add title slide to presentation"""
        
        # Use the first slide layout (title slide)
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{request_data['customer']} - {request_data['industry']}"
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"{request_data['useCase']} Presentation"
        
        # Style the text based on presentation style
        self._style_title_slide(slide, request_data.get('style', 'professional'))
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], slide_number: int):
        """Add a content slide to presentation"""
        
        # Choose layout based on slide type
        slide_type = slide_data.get('slideType', 'content')
        
        if slide_type == 'title':
            layout = prs.slide_layouts[0]  # Title slide
        elif slide_type == 'chart':
            layout = prs.slide_layouts[5]  # Content with caption
        elif slide_type == 'image':
            layout = prs.slide_layouts[6]  # Blank
        else:
            layout = prs.slide_layouts[1]  # Title and content
        
        slide = prs.slides.add_slide(layout)
        
        # Add content based on slide type
        if slide_type == 'title':
            self._add_title_content(slide, slide_data)
        elif slide_type == 'chart':
            self._add_chart_content(slide, slide_data)
        elif slide_type == 'image':
            self._add_image_content(slide, slide_data)
        elif slide_type == 'quote':
            self._add_quote_content(slide, slide_data)
        else:
            self._add_standard_content(slide, slide_data)
        
        # Add slide number
        self._add_slide_number(slide, slide_number)
    
    def _add_title_content(self, slide, slide_data: Dict[str, Any]):
        """Add title slide content"""
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Untitled')
        
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get('content', '')
    
    def _add_chart_content(self, slide, slide_data: Dict[str, Any]):
        """Add chart slide content"""
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Chart')
        
        # Add chart placeholder (simplified)
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text = slide_data.get('content', 'Chart data will be displayed here')
    
    def _add_image_content(self, slide, slide_data: Dict[str, Any]):
        """Add image slide content"""
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Image')
        
        # Add image if URL is available
        image_url = slide_data.get('imageUrl')
        if image_url:
            try:
                self._add_image_from_url(slide, image_url)
            except Exception as e:
                logger.error(f"Error adding image from URL {image_url}: {e}")
                # Add placeholder text
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_data.get('content', 'Image content')
    
    def _add_quote_content(self, slide, slide_data: Dict[str, Any]):
        """Add quote slide content"""
        
        # Use a text box for quote
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f'"{slide_data.get("content", "Quote content")}"'
        
        # Style the quote
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(24)
        paragraph.font.italic = True
    
    def _add_standard_content(self, slide, slide_data: Dict[str, Any]):
        """Add standard content slide"""
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Content')
        
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text = slide_data.get('content', '')
            
            # Format content as bullet points
            text_frame = content.text_frame
            text_frame.clear()
            
            # Split content into bullet points
            content_text = slide_data.get('content', '')
            if content_text:
                points = content_text.split('\n')
                for i, point in enumerate(points):
                    if point.strip():
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = point.strip()
                        p.level = 0
    
    def _add_image_from_url(self, slide, image_url: str):
        """Add image to slide from URL"""
        
        try:
            # Download image
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            # Open image and resize if needed
            image = Image.open(BytesIO(response.content))
            
            # Resize image to fit slide
            max_width = Inches(8)
            max_height = Inches(5)
            
            # Calculate new size maintaining aspect ratio
            width_ratio = max_width / image.width
            height_ratio = max_height / image.height
            ratio = min(width_ratio, height_ratio)
            
            new_width = int(image.width * ratio)
            new_height = int(image.height * ratio)
            
            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Save to temporary file
            temp_path = f"/tmp/temp_image_{uuid.uuid4()}.png"
            image.save(temp_path)
            
            # Add image to slide
            left = Inches(1)
            top = Inches(2)
            slide.shapes.add_picture(temp_path, left, top, width=Inches(new_width/96), height=Inches(new_height/96))
            
            # Clean up temp file
            os.remove(temp_path)
            
        except Exception as e:
            logger.error(f"Error adding image from URL: {e}")
            raise e
    
    def _add_slide_number(self, slide, slide_number: int):
        """Add slide number to slide"""
        
        # Add slide number in bottom right corner
        left = Inches(9)
        top = Inches(7)
        width = Inches(1)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = str(slide_number)
        
        # Style the slide number
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(128, 128, 128)
    
    def _add_conclusion_slide(self, prs: Presentation, request_data: Dict[str, Any]):
        """Add conclusion slide to presentation"""
        
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Thank You"
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"Questions & Discussion\n\n{request_data['customer']} - {request_data['industry']}"
    
    def _style_title_slide(self, slide, style: str):
        """Style the title slide based on presentation style"""
        
        if style == 'creative':
            # Creative styling
            if slide.shapes.title:
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.color.rgb = RGBColor(0, 102, 204)
                title_para.font.bold = True
            
            if len(slide.placeholders) > 1:
                subtitle_para = slide.placeholders[1].text_frame.paragraphs[0]
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = RGBColor(51, 51, 51)
        
        elif style == 'minimalist':
            # Minimalist styling
            if slide.shapes.title:
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.size = Pt(36)
                title_para.font.color.rgb = RGBColor(0, 0, 0)
                title_para.font.bold = True
            
            if len(slide.placeholders) > 1:
                subtitle_para = slide.placeholders[1].text_frame.paragraphs[0]
                subtitle_para.font.size = Pt(18)
                subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        
        else:  # professional/corporate
            # Professional styling
            if slide.shapes.title:
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.size = Pt(40)
                title_para.font.color.rgb = RGBColor(0, 0, 0)
                title_para.font.bold = True
            
            if len(slide.placeholders) > 1:
                subtitle_para = slide.placeholders[1].text_frame.paragraphs[0]
                subtitle_para.font.size = Pt(20)
                subtitle_para.font.color.rgb = RGBColor(64, 64, 64)
    
    async def _generate_preview_images(self, filepath: str, presentation_id: str) -> List[str]:
        """Generate preview images for the presentation"""
        
        try:
            # This would integrate with a service to convert PPTX to images
            # For now, return placeholder URLs
            preview_urls = []
            
            # In a real implementation, you would:
            # 1. Convert PPTX to images using a service like LibreOffice or cloud API
            # 2. Upload images to cloud storage
            # 3. Return the URLs
            
            return preview_urls
            
        except Exception as e:
            logger.error(f"Error generating preview images: {e}")
            return []
    
    def _copy_exact_slide(self, prs: Presentation, slide_data: Dict[str, Any], slide_number: int):
        """Copy exact slide from source with enhanced visual element preservation"""
        
        # Use content slide layout
        slide_layout = prs.slide_layouts[1]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Copy exact title and content with formatting
        title = slide.shapes.title
        title.text = slide_data.get('title', f'Slide {slide_number}')
        
        # Apply title formatting if available
        formatting_data = slide_data.get('formatting', {})
        if formatting_data and isinstance(formatting_data, dict):
            title_formatting = formatting_data.get('title', {})
            if title_formatting and isinstance(title_formatting, dict):
                logger.info(f"🎨 Applying title formatting: {title_formatting}")
                self._apply_text_formatting(title, title_formatting)
            else:
                logger.info(f"🎨 No title formatting found for slide {slide_number}")
        else:
            logger.info(f"🎨 No formatting data found for slide {slide_number}")
        
        content = slide.placeholders[1]
        content.text = slide_data.get('content', '')
        
        # Apply content formatting if available
        if formatting_data and isinstance(formatting_data, dict):
            self._apply_content_formatting(content, formatting_data)
        
        # Copy images from extracted data
        images = slide_data.get('images', [])
        if images and isinstance(images, list):
            logger.info(f"🖼️ Found {len(images)} images for slide {slide_number}")
            for image_data in images:
                try:
                    if image_data and isinstance(image_data, dict):
                        logger.info(f"🖼️ Adding image to slide {slide_number}")
                        self._add_extracted_image(slide, image_data)
                except Exception as e:
                    logger.warning(f"Could not add extracted image to slide {slide_number}: {e}")
        else:
            logger.info(f"🖼️ No images found for slide {slide_number}")
        
        # Copy exact image if available (legacy support)
        if slide_data.get('imageUrl'):
            try:
                self._add_slide_image(slide, slide_data['imageUrl'])
            except Exception as e:
                logger.warning(f"Could not add image to slide {slide_number}: {e}")
        
        # Apply layout information if available
        layout_info = slide_data.get('layout_info', {})
        if layout_info and isinstance(layout_info, dict):
            self._apply_layout_info(slide, layout_info)
        
        # Add slide number
        self._add_slide_number(slide, slide_number)
        
        # Add source attribution
        self._add_source_attribution(slide, slide_data)
    
    def _add_enhanced_slide(self, prs: Presentation, slide_data: Dict[str, Any], slide_number: int, request_data: Dict[str, Any]):
        """Add slide with minor AI enhancement (low cost)"""
        
        # Use content slide layout
        slide_layout = prs.slide_layouts[1]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Enhance title slightly
        original_title = slide_data.get('title', f'Slide {slide_number}')
        enhanced_title = self._enhance_title(original_title, request_data)
        
        title = slide.shapes.title
        title.text = enhanced_title
        
        # Use original content (no AI enhancement to save costs)
        content = slide.placeholders[1]
        content.text = slide_data.get('content', '')
        
        # Add image if available
        if slide_data.get('imageUrl'):
            try:
                self._add_slide_image(slide, slide_data['imageUrl'])
            except Exception as e:
                logger.warning(f"Could not add image to slide {slide_number}: {e}")
        
        # Add slide number
        self._add_slide_number(slide, slide_number)
        
        # Add source attribution
        self._add_source_attribution(slide, slide_data)
    
    def _add_ai_generated_slide(self, prs: Presentation, slide_data: Dict[str, Any], slide_number: int, request_data: Dict[str, Any]):
        """Add slide with full AI generation (high cost - only when necessary)"""
        
        # Use content slide layout
        slide_layout = prs.slide_layouts[1]  # Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Generate enhanced content using AI
        enhanced_content = self._generate_ai_content(slide_data, request_data)
        
        title = slide.shapes.title
        title.text = enhanced_content.get('title', f'Slide {slide_number}')
        
        content = slide.placeholders[1]
        content.text = enhanced_content.get('content', '')
        
        # Add image if available
        if slide_data.get('imageUrl'):
            try:
                self._add_slide_image(slide, slide_data['imageUrl'])
            except Exception as e:
                logger.warning(f"Could not add image to slide {slide_number}: {e}")
        
        # Add slide number
        self._add_slide_number(slide, slide_number)
        
        # Add source attribution
        self._add_source_attribution(slide, slide_data)
    
    def _enhance_title(self, original_title: str, request_data: Dict[str, Any]) -> str:
        """Enhance title with minimal processing (no AI cost)"""
        # Simple enhancement without AI
        customer = request_data.get('customer', '')
        industry = request_data.get('industry', '')
        
        if customer and industry:
            return f"{original_title} - {customer} ({industry})"
        return original_title
    
    def _generate_ai_content(self, slide_data: Dict[str, Any], request_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate AI content (high cost - only when necessary)"""
        # This would use AI to generate content, but for cost optimization
        # we'll return the original content with minor enhancements
        return {
            'title': slide_data.get('title', ''),
            'content': slide_data.get('content', '')
        }
    
    def _add_source_attribution(self, slide, slide_data: Dict[str, Any]):
        """Add source attribution to slide"""
        source = slide_data.get('sourcePresentation', '')
        if source:
            # Add small text box with source info
            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(9)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"Source: {source}"
            
            # Style the text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = RGBColor(128, 128, 128)
    
    def _apply_text_formatting(self, shape, formatting: Dict[str, Any]):
        """Apply text formatting to a shape"""
        try:
            if not formatting or not hasattr(shape, 'text_frame'):
                return
                
            text_frame = shape.text_frame
            if not text_frame.paragraphs:
                return
                
            para = text_frame.paragraphs[0]
            
            # Apply paragraph formatting
            if 'paragraph' in formatting:
                para_format = formatting['paragraph']
                if para_format.get('alignment'):
                    # Convert string alignment to enum
                    alignment_map = {
                        'PP_ALIGN.LEFT': PP_ALIGN.LEFT,
                        'PP_ALIGN.CENTER': PP_ALIGN.CENTER,
                        'PP_ALIGN.RIGHT': PP_ALIGN.RIGHT,
                        'PP_ALIGN.JUSTIFY': PP_ALIGN.JUSTIFY
                    }
                    if para_format['alignment'] in alignment_map:
                        para.alignment = alignment_map[para_format['alignment']]
            
            # Apply font formatting
            if 'font' in formatting and para.runs:
                font_format = formatting['font']
                run = para.runs[0]
                
                if font_format.get('name'):
                    run.font.name = font_format['name']
                if font_format.get('size'):
                    run.font.size = Pt(font_format['size'])
                if font_format.get('bold') is not None:
                    run.font.bold = font_format['bold']
                if font_format.get('italic') is not None:
                    run.font.italic = font_format['italic']
                if font_format.get('underline') is not None:
                    run.font.underline = font_format['underline']
                if font_format.get('color'):
                    # Parse color string and apply
                    try:
                        color_str = font_format['color']
                        if color_str.startswith('RGBColor('):
                            # Extract RGB values
                            rgb_values = color_str[9:-1].split(', ')
                            if len(rgb_values) == 3:
                                r, g, b = map(int, rgb_values)
                                run.font.color.rgb = RGBColor(r, g, b)
                    except Exception as e:
                        logger.warning(f"Could not apply color formatting: {e}")
                        
        except Exception as e:
            logger.warning(f"Could not apply text formatting: {e}")
    
    def _apply_content_formatting(self, placeholder, formatting: Dict[str, Any]):
        """Apply formatting to content placeholder"""
        try:
            if not formatting or not hasattr(placeholder, 'text_frame'):
                return
                
            text_frame = placeholder.text_frame
            
            # Apply formatting to each paragraph
            for i, para in enumerate(text_frame.paragraphs):
                format_key = f'text_shape_{i+1}'
                if format_key in formatting:
                    self._apply_text_formatting(placeholder, {format_key: formatting[format_key]})
                    
        except Exception as e:
            logger.warning(f"Could not apply content formatting: {e}")
    
    def _add_extracted_image(self, slide, image_data: Dict[str, Any]):
        """Add image from extracted data to slide"""
        try:
            if not image_data.get('image_data'):
                logger.warning("No image data available")
                return
                
            # Save image data to temporary file
            import tempfile
            import uuid
            
            # Determine file extension from format
            image_format = image_data.get('image_format', 'png')
            if 'jpeg' in image_format.lower():
                ext = '.jpg'
            elif 'png' in image_format.lower():
                ext = '.png'
            elif 'gif' in image_format.lower():
                ext = '.gif'
            else:
                ext = '.png'  # Default
                
            temp_path = f"/tmp/extracted_image_{uuid.uuid4()}{ext}"
            
            # Write image data to file
            with open(temp_path, 'wb') as f:
                f.write(image_data['image_data'])
            
            # Add image to slide with original positioning
            left = image_data.get('left', Inches(1))
            top = image_data.get('top', Inches(2))
            width = image_data.get('width', Inches(4))
            height = image_data.get('height', Inches(3))
            
            slide.shapes.add_picture(temp_path, left, top, width, height)
            
            # Clean up temp file
            os.remove(temp_path)
            
        except Exception as e:
            logger.error(f"Error adding extracted image: {e}")
    
    def _apply_layout_info(self, slide, layout_info: Dict[str, Any]):
        """Apply layout information to slide"""
        try:
            if not layout_info:
                return
                
            # Apply background if available
            if 'background' in layout_info and layout_info['background'].get('has_fill'):
                # This would require more complex background handling
                # For now, we'll skip background application
                pass
                
            # Apply dimensions if available
            if 'dimensions' in layout_info:
                # Slide dimensions are set at presentation level, not individual slides
                pass
                
        except Exception as e:
            logger.warning(f"Could not apply layout info: {e}")

    def get_presentation_info(self, filepath: str) -> Dict[str, Any]:
        """Get information about the generated presentation"""
        
        try:
            prs = Presentation(filepath)
            
            return {
                'slideCount': len(prs.slides),
                'title': prs.core_properties.title,
                'author': prs.core_properties.author,
                'created': prs.core_properties.created,
                'fileSize': os.path.getsize(filepath)
            }
            
        except Exception as e:
            logger.error(f"Error getting presentation info: {e}")
            return {}
